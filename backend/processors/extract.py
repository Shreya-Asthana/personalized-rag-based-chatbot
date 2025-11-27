import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
import os

def extract_text(path):
    ext = path.split('.')[-1].lower()

    if ext == "pdf":
        with pdfplumber.open(path) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)

    if ext == "docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == "pptx":
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    if ext in ["xlsx", "xls"]:
        df = pd.read_excel(path)
        return df.to_string()

    if ext == "txt":
        return open(path).read()

    return ""